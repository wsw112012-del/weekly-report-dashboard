"""
주간보고 PPT 생성 스크립트
사용법: python make_report.py 데이터 [PDF파일명 또는 전체경로]
        python make_report.py 페이먼트 [PDF파일명 또는 전체경로]
       PDF파일명만 넘기면 Desktop · 다운로드 폴더에서 자동 탐색
"""

import os
import re
import sys
import shutil
import tempfile
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Mm
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

FONT_NAME = '맑은 고딕'
CM = 360000  # 1cm in EMU

# A4 가로 슬라이드 크기
SLIDE_W_IN = 297 / 25.4   # ≈ 11.693"
SLIDE_H_IN = 210 / 25.4   # ≈  8.268"

BASE_PATH = r"C:\Users\쿠콘_우승우\Desktop\업무\00. '26 쿠콘전략실\08. 주간보고\데이터전략센터\주간보고"

# 표제 고정 문구
HEADERS = {
    '데이터': {
        'title':    '1) 데이터 정책 동향',
        'subtitle': '- 데이터3법 및 유관기관(금융위, 개보위 등) 정책동향 공유',
    },
    '페이먼트': {
        'title':    '2) 페이먼트 정책 동향',
        'subtitle': '- 전자금융업법 및 특정금융정보법 유관기관(금융위, 금감원 등) 정책동향 공유',
    },
    'AML': {
        'title':    '3) AML · 제재 정책 동향',
        'subtitle': '- 자금세탁방지법 및 특정금융정보법 유관기관 정책동향 공유',
    },
}

# 본문 들여쓰기: (marL cm, 첫줄 내어쓰기 cm)
INDENT = {
    'diamond': (1.01, -0.48),
    'summary': (1.48, -0.48),
    'section': (2.00, -0.48),
    'detail':  (2.48, -0.48),
}

SPACE_BEFORE = {
    'diamond': {'first': 10, 'rest': 14},
    'summary': 0,
    'section': 0,
    'detail':  1,
}


def detect_type(stripped: str, level: int) -> str:
    if stripped.startswith('◆'):
        return 'diamond'
    if stripped.startswith('•'):
        return 'section'
    if stripped.startswith('-'):
        return 'summary' if level == 1 else 'detail'
    return 'detail'


def apply_line_spacing(p, multiple: float = 1.5):
    pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.find(qn('a:lnSpc'))
    if lnSpc is None:
        lnSpc = etree.Element(qn('a:lnSpc'))
        pPr.insert(0, lnSpc)
    else:
        for child in list(lnSpc):
            lnSpc.remove(child)
    spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
    spcPct.attrib['val'] = str(int(multiple * 100000))


def apply_indent(p, line_type: str):
    mar_cm, ind_cm = INDENT[line_type]
    pPr = p._p.get_or_add_pPr()
    pPr.attrib['marL'] = str(int(mar_cm * CM))
    pPr.attrib['indent'] = str(int(ind_cm * CM))


def apply_font(run, size_pt: float, bold: bool = False):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    rPr.attrib['spc'] = '0'
    for tag in [qn('a:latin'), qn('a:ea'), qn('a:cs')]:
        elem = rPr.find(tag)
        if elem is None:
            elem = etree.SubElement(rPr, tag)
        elem.attrib['typeface'] = FONT_NAME


def add_header(tf, report_type: str):
    hdr = HEADERS[report_type]

    p = tf.paragraphs[0]
    apply_line_spacing(p, 1.5)
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = hdr['title']
    apply_font(run, size_pt=14, bold=True)

    p2 = tf.add_paragraph()
    apply_line_spacing(p2, 1.5)
    p2.space_before = Pt(0)
    p2.space_after = Pt(0)
    p2.alignment = PP_ALIGN.LEFT
    pPr = p2._p.get_or_add_pPr()
    pPr.attrib['marL'] = str(int(0.5 * CM))
    pPr.attrib['indent'] = '0'
    run2 = p2.add_run()
    run2.text = hdr['subtitle']
    apply_font(run2, size_pt=13, bold=False)


def extract_press_date(content: str) -> str:
    for line in content.split('\n'):
        if line.strip().startswith('◆'):
            m = re.search(r"'(\d{2})\.\s*(\d{1,2})\.\s*(\d{1,2})", line)
            if m:
                yy, mm, dd = m.groups()
                return f"20{yy}{int(mm):02d}{int(dd):02d}"
    return date.today().strftime('%Y%m%d')


def extract_press_title(content: str) -> str:
    """◆ 줄에서 「」 안의 핵심 제목 추출 — 파일명용"""
    for line in content.strip().split('\n'):
        if line.strip().startswith('◆'):
            m = re.search(r'「(.+?)」', line)
            if m:
                name = m.group(1).strip()
                name = re.sub(r'[\\/:*?"<>|]', '', name)
                return name[:40]
    return ""


def build_output_path(content: str, report_type: str) -> str:
    today = date.today().strftime('%Y%m%d')
    press_date = extract_press_date(content)
    press_title = extract_press_title(content)
    folder = os.path.join(BASE_PATH, today)
    os.makedirs(folder, exist_ok=True)
    base = press_title if press_title else report_type
    return os.path.join(folder, f"{base}_{press_date}.pptx")


def create_pptx(content: str, output_path: str, report_type: str):
    prs = Presentation()
    prs.slide_width  = Mm(297)   # A4 가로
    prs.slide_height = Mm(210)   # A4 세로
    # type="screen4x3" is the default; must be changed to "custom" so
    # PowerPoint respects our cx/cy values instead of the preset dimensions
    prs_el = prs._element
    sld_sz = prs_el.find(qn('p:sldSz'))
    if sld_sz is not None:
        sld_sz.attrib['type'] = 'custom'

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(
        Mm(10), Mm(12),
        Mm(275), Mm(180),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    add_header(tf, report_type)

    diamond_count = 0
    for line in content.strip().split('\n'):
        stripped = line.strip()
        if not stripped:
            continue

        level = (len(line) - len(line.lstrip(' '))) // 2
        line_type = detect_type(stripped, level)

        p = tf.add_paragraph()
        apply_line_spacing(p, 1.5)

        if line_type == 'diamond':
            p.space_before = Pt(SPACE_BEFORE['diamond']['first'] if diamond_count == 0 else SPACE_BEFORE['diamond']['rest'])
            diamond_count += 1
        else:
            p.space_before = Pt(SPACE_BEFORE[line_type])
        p.space_after = Pt(0)

        p.alignment = PP_ALIGN.LEFT
        apply_indent(p, line_type)

        run = p.add_run()
        run.text = stripped
        apply_font(run, size_pt=10.5 if line_type == 'diamond' else 10,
                   bold=(line_type in ('diamond', 'section')))

    prs.save(output_path)
    print(f"저장 완료: {output_path}")


def find_pdf(name: str) -> str | None:
    if os.path.isabs(name) or os.sep in name or '/' in name:
        return name
    home = os.path.expanduser('~')
    for d in [os.path.join(home, 'Desktop'),
              os.path.join(home, 'Downloads'),
              os.path.join(home, '다운로드')]:
        candidate = os.path.join(d, name)
        if os.path.exists(candidate):
            return candidate
    return None


def make_short_label(pdf_name: str, max_len: int = 25) -> str:
    base = os.path.splitext(os.path.basename(pdf_name))[0]
    if len(base) <= max_len:
        return base
    return base[:max_len - 3] + '...'


def insert_pdf_object(pptx_path: str, pdf_path: str):
    import win32com.client

    pdf_path  = os.path.abspath(pdf_path)
    pptx_path = os.path.abspath(pptx_path)
    label     = make_short_label(pdf_path)
    tmp_dir   = tempfile.mkdtemp()
    short_pdf = os.path.join(tmp_dir, label + '.pdf')

    try:
        shutil.copy2(pdf_path, short_pdf)

        ppt = win32com.client.Dispatch("PowerPoint.Application")
        ppt.Visible = True
        try:
            prs   = ppt.Presentations.Open(pptx_path)
            slide = prs.Slides(1)

            # A4 가로 기준 (pt = 1/72")
            slide_w_pt = SLIDE_W_IN * 72   # ≈ 841.9 pt
            slide_h_pt = SLIDE_H_IN * 72   # ≈ 595.3 pt
            obj_w  = 2.5 * 72
            obj_h  = 1.8 * 72
            margin = 0.15 * 72
            obj_left = slide_w_pt - obj_w - margin
            obj_top  = slide_h_pt - obj_h - margin

            slide.Shapes.AddOLEObject(
                Left=obj_left, Top=obj_top,
                Width=obj_w,   Height=obj_h,
                FileName=short_pdf,
                Link=False,
                DisplayAsIcon=True,
            )

            prs.Save()
            prs.Close()
            print(f"PDF 개체 삽입 완료: {label}.pdf")
        finally:
            ppt.Quit()
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


# ── 여기에 보고 내용을 붙여넣으세요 ──────────────────────────────────
CONTENT = """
◆ 금융위원회 | 「가상자산시장 시세조종 혐의자 수사기관 통보 - 제8차 금융위원회 정례회의」  '26.4.29(수)
  - 가상자산시장 시세조종 혐의자 수사기관 통보 - 제8차 금융위원회 정례회의('26.4.29.

    • 주요 내용
      - ... 서비스 접속 차단 , API 서비스 제한 등의 조치를 받을 수 있다
      - 서비스 접속 차단 , API 서비스 제한 등의 조치를 받을 수 있다. 나아가, 타인에게 대여한 본인의 API 키 (Key) 가 불공정거래,자금세탁등에 사용 될 경우 이용자 (명의 자) 는 불법행위 의 공범 으로 처벌받는 등 민
      - 형사상 책임 을 질 수 있으 므로 각별히

◆ 금융위원회 | 「대부업법」  '26.4.28(화)
  - 불법사금융 신고의 문턱은 낮추고 — 범죄 차단속도는 높이겠습니다. - 「대부업법」 시행령 개정

    • 주요 내용
      - 채무조정을 신청하여 변제 중이었으나, 이후 영업 부진으로 직원 급여 등 운영자금마련에 어려움이 발생함
      - 이에 '25.6월 인터넷 대부중개 사이트를 통해 불법사금융을   59건 의 의심계좌 를 해당 금융회사에 통보 하였다
      - 금융회사는자금세탁방지제도에 따라 계좌 명의인에 대해 거래자금의 원천, 금융거래의 목적 등을 소명하도록
"""
# ───────────────────────────────────────────────────────────────────────

if __name__ == '__main__':
    if len(sys.argv) < 2 or sys.argv[1] not in ('데이터', '페이먼트', 'AML'):
        print("사용법: python make_report.py 데이터 [PDF파일명]")
        print("        python make_report.py 페이먼트 [PDF파일명]")
        print("        python make_report.py AML [PDF파일명]")
        sys.exit(1)

    report_type = sys.argv[1]
    pdf_arg     = sys.argv[2] if len(sys.argv) >= 3 else None

    out = build_output_path(CONTENT, report_type)
    create_pptx(CONTENT, out, report_type)

    if pdf_arg:
        resolved = find_pdf(pdf_arg)
        if resolved:
            insert_pdf_object(out, resolved)
        else:
            print(f"PDF 파일을 찾을 수 없습니다: {pdf_arg}")
